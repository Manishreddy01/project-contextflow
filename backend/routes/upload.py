from fastapi import APIRouter, UploadFile, File, Form, HTTPException
from typing import List, Optional
import os
import sys
from pathlib import Path


from rag.vectorstore.store_vectors import store_vectors

# PDF + PPTX extraction
from pdfminer.high_level import extract_text as pdf_extract_text
from pptx import Presentation

router = APIRouter()
print(">>> USING UPLOAD.PY AT PATH:", __file__)
print("🔥🔥🔥 ACTIVE UPLOAD ROUTE:", __file__)



def extract_txt(content: bytes) -> str:
    return content.decode("utf-8", errors="ignore")


def extract_pdf(file_path: str) -> str:
    try:
        return pdf_extract_text(file_path)
    except:
        return ""


def extract_pptx(file_path: str) -> str:
    try:
        prs = Presentation(file_path)
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(shape.text)
        return "\n".join(text_runs)
    except:
        return ""


@router.post("/")
async def upload_files(
    files: Optional[List[UploadFile]] = File(None),
    text: Optional[str] = Form(None),
    conversationId: Optional[str] = Form(None)
):
    if not conversationId:
        raise HTTPException(status_code=400, detail="Missing conversationId")

    if not files and not text:
        return {"success": False, "message": "No content received"}

    try:
        uploaded_names = []

        # Handle uploaded files
        if files:
            for file in files:
                # Fix: sanitize filename (max 50 chars)
                print("🔥 ENTERED FILE LOOP — REAL CODE RUNNING")
                orig = file.filename
                ext = Path(orig).suffix.lower()
                if not ext:
                    ext = ".txt"
                if len(orig) > 50:
                    safe_name = f"uploaded{ext}"
                else:
                    safe_name = orig

                # Log for debugging
                print(">>> SAFE FILENAME:", safe_name)

                uploaded_names.append(safe_name)
                temp_path = f"/tmp/{safe_name}"
                with open(temp_path, "wb") as f:
                    f.write(await file.read())

                # Extract based on type
                if safe_name.endswith(".txt"):
                    with open(temp_path, "rb") as f:
                        extracted = extract_txt(f.read())
                elif safe_name.endswith(".pdf"):
                    extracted = extract_pdf(temp_path)
                elif safe_name.endswith(".pptx"):
                    extracted = extract_pptx(temp_path)
                else:
                    extracted = ""

                if extracted.strip():
                    store_vectors(
                        content=extracted,
                        file_name=safe_name,
                        conversation_id=conversationId
                    )

        # Handle pasted text
       # if text:
            #store_vectors(
                #content=text,
                #file_name="pasted_text.txt",
               # conversation_id=conversationId
           # )
        print("UPLOAD conversationId:", conversationId)
        return {
            "success": True,
            "uploadedFiles": uploaded_names,
            "textReceived": bool(text)
        }

    except Exception as e:
        print(f"[ERROR] Upload failed: {e}")
        raise HTTPException(status_code=500, detail="Upload failed. Check logs.")
