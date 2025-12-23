from typing import List, Union
from langchain_core.documents import Document
from langchain_community.document_loaders import (
    PyPDFLoader,
    TextLoader,
    Docx2txtLoader,
)
from pathlib import Path
import tempfile

# For PPTX support
from pptx import Presentation


def extract_pptx(file_path: str) -> str:
    """Extract text from a .pptx file."""
    prs = Presentation(file_path)
    text_runs = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                text_runs.append(shape.text)

    return "\n".join(text_runs)


def load_documents(file_content_or_path: Union[str, bytes], file_name: str) -> List[Document]:
    """
    Unified document loader:
    - Accepts bytes (uploaded files)
    - Accepts file path (local files)
    - Accepts raw text as string
    - Supports TXT, PDF, DOCX, PPTX
    """
    suffix = Path(file_name).suffix.lower()

    # ------------------------------------------------------------
    # Case 1: file_content_or_path is an actual PATH
    # ------------------------------------------------------------
    if isinstance(file_content_or_path, str) and Path(file_content_or_path).exists():
        file_path = file_content_or_path

    # ------------------------------------------------------------
    # Case 2: content is raw text in string form -> write to temp
    # ------------------------------------------------------------
    else:
        if isinstance(file_content_or_path, str):
            file_content_or_path = file_content_or_path.encode("utf-8")

        with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
            tmp.write(file_content_or_path)
            tmp.flush()
            file_path = tmp.name

    # ------------------------------------------------------------
    # Loader selection based on file extension
    # ------------------------------------------------------------
    if suffix == ".pdf":
        loader = PyPDFLoader(file_path)
        documents = loader.load()

    elif suffix == ".txt":
        loader = TextLoader(file_path)
        documents = loader.load()

    elif suffix == ".docx":
        try:
            from docx import Document as DocxReader
            doc = DocxReader(file_path)
            text_runs = [p.text for p in doc.paragraphs if p.text.strip()]
            content = "\n".join(text_runs)
            documents = [Document(page_content=content, metadata={"source": file_name})]
        except Exception as e:
            print("❌ DOCX extraction failed:", e)
            return []
    elif suffix == ".pptx":
        extracted = extract_pptx(file_path)
        documents = [Document(page_content=extracted, metadata={"source": file_name})]

    else:
        raise ValueError(f"❌ Unsupported file type: {suffix}")

    print(f"📄 Loaded {len(documents)} pages from {file_name}")
    return documents
